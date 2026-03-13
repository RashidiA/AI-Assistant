import streamlit as st
from streamlit_webrtc import webrtc_streamer, WebRtcMode
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import os
import glob
import queue
import pydub
from io import BytesIO
from duckduckgo_search import DDGS
import speech_recognition as sr
from datetime import datetime

# --- CONFIG ---
st.set_page_config(page_title="AI Engineering Assistant", layout="wide")

# Initialize session state for commands
if 'last_cmd' not in st.session_state:
    st.session_state['last_cmd'] = ""

# --- 1. VOICE PROCESSING (BROWSER COMPATIBLE) ---
def process_audio(audio_frames):
    """Converts raw browser audio frames into text."""
    if not audio_frames:
        return "No audio captured."
    
    # Merge chunks into one audio segment
    sound = pydub.AudioSegment.empty()
    for frame in audio_frames:
        s = pydub.AudioSegment(
            data=frame.to_ndarray().tobytes(),
            sample_width=2, # Standard 16-bit
            frame_rate=frame.sample_rate,
            channels=1
        )
        sound += s
    
    # Export to memory
    buffer = BytesIO()
    sound.export(buffer, format="wav")
    buffer.seek(0)
    
    r = sr.Recognizer()
    with sr.AudioFile(buffer) as source:
        audio_data = r.record(source)
        try:
            return r.recognize_google(audio_data)
        except Exception:
            return "Could not understand audio."

# --- 2. ENGINEERING TOOLS ---
def analyze_csv(file_path):
    """Simple Cpk Analysis for Geely Engineering."""
    try:
        df = pd.read_csv(file_path)
        numeric_data = df.select_dtypes(include=[np.number])
        if numeric_data.empty:
            return None, None
        
        data = numeric_data.iloc[:, 0]
        mean = data.mean()
        std = data.std()
        # Example LSL/USL: 9.5 to 10.5
        cpk = min((10.5 - mean)/(3*std), (mean - 9.5)/(3*std)) if std > 0 else 0
        return mean, cpk
    except Exception:
        return None, None

def web_research(query):
    """Search DuckDuckGo for engineering standards."""
    try:
        with DDGS() as ddgs:
            return [r for r in ddgs.text(f"engineering standard {query}", max_results=3)]
    except Exception:
        return []

# --- 3. UI LAYOUT ---
st.title("🤖 Pro AI Engineering Assistant")
st.write("Voice or type commands to analyze data, search standards, or make reports.")

col1, col2 = st.columns([1, 2])

with col1:
    st.subheader("🎙️ Voice Control")
    webrtc_ctx = webrtc_streamer(
        key="speech-to-text",
        mode=WebRtcMode.SENDONLY,
        audio_receiver_size=1024,
        media_stream_constraints={"video": False, "audio": True},
    )

    if webrtc_ctx.state.playing:
        st.info("Listening... Click 'Stop' when done speaking.")
        # Frames are gathered in the background by the audio_receiver

    if not webrtc_ctx.state.playing and webrtc_ctx.audio_receiver:
        # User just clicked stop, let's process
        with st.spinner("Processing voice..."):
            try:
                frames = webrtc_ctx.audio_receiver.get_frames(timeout=1)
                res = process_audio(frames)
                st.session_state['last_cmd'] = res
            except Exception:
                pass

with col2:
    st.subheader("⌨️ Command Interface")
    user_input = st.text_input("Current Command:", value=st.session_state['last_cmd'])
    
    if user_input:
        cmd = user_input.lower()
        
        # LOGIC 1: DATA ANALYSIS & PPT
        if "analyze" in cmd or "ppt" in cmd:
            csv_files = glob.glob("*.csv")
            if csv_files:
                target = csv_files[0]
                mean, cpk = analyze_csv(target)
                
                if mean is not None:
                    st.success(f"File: {target} | Mean: {mean:.3f} | Cpk: {cpk:.2f}")
                    
                    # Generate PPT
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f"Engineering Report: {target}"
                    slide.placeholders[1].text = f"Date: {datetime.now().strftime('%Y-%m-%d')}\nMean: {mean:.4f}\nCpk: {cpk:.2f}"
                    
                    ppt_io = BytesIO()
                    prs.save(ppt_io)
                    st.download_button("📥 Download PowerPoint", ppt_io.getvalue(), "Geely_Report.pptx")
                else:
                    st.error("CSV file contains no numeric data.")
            else:
                st.warning("No CSV files found in directory. Please upload one to GitHub.")

        # LOGIC 2: WEB SEARCH
        elif "search" in cmd:
            query = cmd.replace("search", "").strip()
            with st.spinner(f"Searching for {query}..."):
                results = web_research(query)
                for r in results:
                    with st.expander(r['title']):
                        st.write(r['body'])
                        st.caption(f"Source: {r['href']}")

# --- SIDEBAR HELP ---
st.sidebar.header("How to Use")
st.sidebar.markdown("""
1. **Upload a CSV** to your GitHub repo.
2. **Click Start** on the Mic and say: *"Analyze data and make a ppt"*
3. **Or type**: *"Search ISO 9001"*
""")
